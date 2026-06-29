"""
Presentación de Defensa (PowerPoint) — EFT "Banco Digital Chile"
=================================================================
Genera un deck ejecutivo (16:9) para la defensa individual de 10 minutos, con
notas del orador en cada diapositiva y los gráficos reales del dashboard.

Paleta "Midnight Executive" (banca): navy dominante + ámbar de acento.
Uso:
    python run_observability.py            # genera gráficos/KPIs
    python docs/build_presentacion.py
"""

from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from observability import metrics as M

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CHARTS = os.path.join(ROOT, "dashboard", "charts")
ENTREGA = os.path.join(ROOT, "entrega"); os.makedirs(ENTREGA, exist_ok=True)
SALIDA = os.path.join(ENTREGA, "EFT_Presentacion_Defensa.pptx")

# Paleta
NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x0F, 0x25, 0x40)
ICE = RGBColor(0xD6, 0xE4, 0xF0)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LGREY = RGBColor(0xB8, 0xC4, 0xD0)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC0, 0x39, 0x2B)
DARKTXT = RGBColor(0x22, 0x2B, 0x35)

HFONT = "Georgia"
BFONT = "Calibri"

W, H = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    _solid(s, color)
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1.25); s.fill.background()
    s.shadow.inherit = False
    return s


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=2, line_spacing=None):
    """runs: lista de párrafos; cada párrafo es lista de (txt,size,bold,color,name?,italic?)."""
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, parr in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        if line_spacing: p.line_spacing = line_spacing
        for seg in parr:
            txt, size, bold, color = seg[0], seg[1], seg[2], seg[3]
            name = seg[4] if len(seg) > 4 else BFONT
            italic = seg[5] if len(seg) > 5 else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = name; r.font.italic = italic
    return tb


def bullets(slide, l, t, w, h, items, size=15, color=DARKTXT, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        # viñeta ámbar
        rb = p.add_run(); rb.text = "▪  "; rb.font.size = Pt(size)
        rb.font.color.rgb = AMBER; rb.font.bold = True; rb.font.name = BFONT
        if isinstance(it, tuple):
            rt = p.add_run(); rt.text = it[0]; rt.font.bold = True
            rt.font.size = Pt(size); rt.font.color.rgb = color; rt.font.name = BFONT
            rr = p.add_run(); rr.text = it[1]; rr.font.size = Pt(size)
            rr.font.color.rgb = color; rr.font.name = BFONT
        else:
            r = p.add_run(); r.text = it; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.name = BFONT
    return tb


def notas(slide, texto):
    slide.notes_slide.notes_text_frame.text = texto


def kicker(slide, num, titulo):
    """Encabezado de sección: número en cuadro ámbar + título (sin línea de acento)."""
    rect(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5), AMBER,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.5),
         [[(num, 18, True, DARK, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, Inches(1.25), Inches(0.5), Inches(11.4), Inches(0.7),
         [[(titulo, 30, True, NAVY, HFONT)]], anchor=MSO_ANCHOR.MIDDLE)


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    df = M.load_df()
    g = M.global_metrics(df)
    sec = M.seguridad(df)

    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    # ── 1. Portada (dark) ─────────────────────────────────────────────────────
    s = slide_blank(prs)
    rect(s, 0, 0, W, H, DARK)
    rect(s, Inches(0.6), Inches(2.55), Inches(0.9), Inches(0.12), AMBER)
    text(s, Inches(0.6), Inches(2.75), Inches(12), Inches(2),
         [[("Banco Digital Chile", 46, True, WHITE, HFONT)],
          [("Solución de IA conversacional para banca digital", 22, False, ICE, BFONT)]])
    text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(1.2),
         [[("Defensa — Evaluación Final Transversal · ISY0101", 15, True, AMBER, BFONT)],
          [("Roger Rosas · Rodrigo Santis · Edgardo Gutiérrez", 14, False, LGREY, BFONT)],
          [("Junio 2025", 12, False, LGREY, BFONT)]])
    text(s, Inches(9.4), Inches(0.6), Inches(3.3), Inches(0.5),
         [[("REUNIÓN EJECUTIVA", 12, True, LGREY, BFONT)]], align=PP_ALIGN.RIGHT)
    notas(s, "Saludo. Preséntate y abre el contexto: somos un equipo que diseñó un agente de IA "
             "para Banco Digital Chile. En 10 minutos mostraré el caso, la solución, una demo, "
             "los resultados de observabilidad y las mejoras propuestas. (~30 s)")

    # ── 2. El caso ────────────────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "1", "El caso organizacional")
    bullets(s, Inches(0.6), Inches(1.7), Inches(6.6), Inches(4.8), [
        ("Banco 100% digital. ", "Alto volumen de consultas repetitivas y operaciones sensibles."),
        ("Atención humana no escala. ", "Costos altos y tiempos de espera."),
        ("Chatbot de reglas insuficiente. ", "No entiende lenguaje natural ni ejecuta acciones."),
        ("Exigencias. ", "Trazabilidad, seguridad y cumplimiento (Ley 19.628, CMF)."),
    ], size=17, gap=14)
    # callout
    rect(s, Inches(7.6), Inches(1.9), Inches(5.1), Inches(3.9), ICE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(7.9), Inches(2.2), Inches(4.5), Inches(3.4),
         [[("El desafío", 18, True, NAVY, HFONT)],
          [("", 8, False, NAVY)],
          [("Automatizar la atención de extremo a extremo", 16, True, DARKTXT, BFONT)],
          [("manteniendo contexto, seguridad y trazabilidad, "
            "sin sacrificar la experiencia del cliente.", 14, False, GREY, BFONT)]],
         anchor=MSO_ANCHOR.TOP)
    notas(s, "Explica el problema y por qué importa: volumen, costo, limitaciones del chatbot "
             "tradicional. Cierra con el desafío del recuadro. (~1 min)")

    # ── 3. La solución ────────────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "2", "La solución propuesta")
    text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
         [[("Un agente conversacional que consulta, opera y razona — con LLM + RAG, memoria, "
            "planificación y orquestación multi-agente.", 15, False, GREY, BFONT, True)]])
    cards = [
        ("Consulta", "Saldos, productos\ny condiciones"),
        ("Opera", "Transferencias\ny reclamos"),
        ("Razona", "Simulación de\ncréditos"),
        ("Recuerda", "Memoria CP/LP\npor cliente"),
        ("Coordina", "Orquestación\nmulti-agente"),
    ]
    x = Inches(0.6); cw = Inches(2.32); gap = Inches(0.12)
    for titulo, sub in cards:
        rect(s, x, Inches(2.5), cw, Inches(2.7), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, Inches(2.85), cw, Inches(2), [
            [(titulo, 19, True, WHITE, HFONT)],
            [("", 6, False, WHITE)],
            *[[(linea, 13, False, ICE, BFONT)] for linea in sub.split("\n")],
        ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        x = Emu(x + cw + gap)
    notas(s, "Describe qué hace el agente: 5 capacidades. Enfatiza que NO es solo un chatbot: "
             "ejecuta acciones reales y coordina especialistas. (~1 min)")

    # ── 4. Arquitectura ───────────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "3", "Arquitectura de la solución")
    s.shapes.add_picture(os.path.join(CHARTS, "00_arquitectura.png"),
                         Inches(2.2), Inches(1.45), height=Inches(5.7))
    notas(s, "Recorre las capas de arriba abajo: seguridad de entrada → agente LangChain (ReAct) "
             "+ LLM → herramientas/RAG/memoria/orquestación → observabilidad de salida. "
             "Prepárate para '¿por qué LangChain y CrewAI?'. (~1 min)")

    # ── 5. Agente funcional + demo ────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "4", "Agente funcional  ·  demostración")
    cols = [
        ("Herramientas", ["consultar_saldo", "realizar_transferencia",
                           "calcular_cuota_credito", "registrar_reclamo",
                           "buscar_producto (RAG)"]),
        ("Memoria", ["Corto plazo: ventana k=5", "Largo plazo: JSON por RUT",
                     "Continuidad entre sesiones"]),
        ("Planificación", ["Plan-and-Execute", "CrewAI: 4 roles", "Delegación jerárquica"]),
    ]
    x = Inches(0.6); cw = Inches(3.95); gap = Inches(0.2)
    for titulo, its in cols:
        rect(s, x, Inches(1.7), cw, Inches(0.6), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, x, Inches(1.7), cw, Inches(0.6), [[(titulo, 16, True, WHITE, HFONT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bullets(s, Emu(x + Inches(0.15)), Inches(2.55), Emu(cw - Inches(0.3)), Inches(3),
                its, size=13.5, gap=8)
        x = Emu(x + cw + gap)
    rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.95), ICE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.95),
         [[("DEMO EN VIVO:  ", 14, True, NAVY, BFONT),
           ("consulta de saldo con RUT → simulación de crédito multi-etapa → dashboard de métricas.",
            14, False, DARKTXT, BFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    notas(s, "Aquí ejecutas la DEMO real (ten el entorno listo). Muestra memoria (recuerda el "
             "saldo) y planificación (crédito multi-etapa). Si no hay API, usa el notebook. (~2 min)")

    # ── 6. LLM + RAG ──────────────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "5", "Diseño LLM + RAG")
    bullets(s, Inches(0.6), Inches(1.7), Inches(6.5), Inches(4.5), [
        ("Prompts. ", "System prompt con rol y reglas + few-shot + chain-of-thought."),
        ("RAG. ", "Chunking → TF-IDF → recuperación top-k para fundamentar respuestas."),
        ("Fuentes internas. ", "Productos y políticas del banco."),
        ("Fuentes externas. ", "Normativa CMF, SERNAC y Ley 19.628."),
    ], size=16, gap=12)
    rect(s, Inches(7.5), Inches(1.9), Inches(5.2), Inches(4.2), ICE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(7.8), Inches(2.15), Inches(4.6), Inches(3.8), [
        [("Recuperación combinada", 16, True, NAVY, HFONT)],
        [("Consulta: tasa de crédito + deber de información", 12, False, GREY, BFONT, True)],
        [("", 6, False, GREY)],
        [("● interna — Catálogo de productos", 13, False, DARKTXT, BFONT)],
        [("● externa — CMF (información al cliente)", 13, False, DARKTXT, BFONT)],
        [("● interna — Política de crédito", 13, False, DARKTXT, BFONT)],
        [("", 6, False, GREY)],
        [("→ combina ambos orígenes en una sola respuesta", 12, True, GREEN, BFONT)],
    ])
    notas(s, "Explica que el RAG combina fuentes internas y externas (clave de la rúbrica). "
             "El ejemplo muestra recuperación de ambos orígenes. Menciona que en producción se "
             "usan embeddings densos. (~1 min)")

    # ── 7. Observabilidad: KPIs + dashboard ───────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "6", "Observabilidad — resultados")
    kpis = [(f"{g['precision_exito']*100:.0f}%", "Precisión", GREEN),
            (f"{g['tasa_error']*100:.1f}%", "Tasa de error", AMBER),
            (f"{g['latencia_p95_ms']/1000:.1f}s", "Latencia p95", NAVY),
            (f"{sec['tasa_bloqueo']*100:.0f}%", "Bloqueo injection", RED)]
    x = Inches(0.6); cw = Inches(2.9); gap = Inches(0.13)
    for val, lab, col in kpis:
        rect(s, x, Inches(1.65), cw, Inches(1.5), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             line=col)
        text(s, x, Inches(1.72), cw, Inches(1.4),
             [[(val, 34, True, col, HFONT)], [(lab, 12, False, GREY, BFONT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(x + cw + gap)
    s.shapes.add_picture(os.path.join(CHARTS, "01_precision_por_escenario.png"),
                         Inches(0.6), Inches(3.5), width=Inches(6.0))
    s.shapes.add_picture(os.path.join(CHARTS, "02_latencia_por_escenario.png"),
                         Inches(6.8), Inches(3.5), width=Inches(6.0))
    notas(s, f"Presenta las métricas reales sobre {g['interacciones_totales']} interacciones: "
             "precisión 85%, error 6.9%, p95 8.3s, bloqueo de injection 88%. El panel es el "
             "dashboard de monitoreo. (~1.5 min)")

    # ── 8. Trazabilidad: hallazgos ────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "7", "Trazabilidad — hallazgos clave")
    s.shapes.add_picture(os.path.join(CHARTS, "03_serie_temporal_incidente.png"),
                         Inches(6.7), Inches(1.7), width=Inches(6.2))
    bullets(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5), [
        ("Incidente detectado. ", "2025-06-11, latencia y error con z≈3 (auto-detección)."),
        ("Cuello de botella. ", "Orquestación multi-agente, p95 ≈ 21 s."),
        ("Predictor de éxito. ", "Elegir la herramienta correcta: 95% vs 0%."),
        ("Ambigüedad. ", "Precisión cae 37 pp con entradas poco claras."),
    ], size=15.5, gap=13)
    notas(s, "Muestra que la observabilidad SIRVE: detectó un incidente automáticamente, "
             "identificó el cuello de botella y la causa raíz de los fallos (selección de "
             "herramienta y ambigüedad). Demuestra capacidad de análisis. (~1 min)")

    # ── 9. Seguridad y ética ──────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "8", "Seguridad y uso responsable")
    bullets(s, Inches(0.6), Inches(1.75), Inches(6.4), Inches(5), [
        ("Privacidad por diseño. ", "Enmascaramiento de PII (RUT, tarjeta) en logs y prompts."),
        ("Anti prompt-injection. ", "Bloqueo de manipulación del prompt (OWASP LLM01)."),
        ("Rate-limit + validación. ", "Control de abuso y saneamiento de entrada."),
        ("Auditoría con hash. ", "Bitácora con integridad (no-repudio)."),
        ("Confirmación humana. ", "Toda transferencia se valida explícitamente."),
    ], size=15.5, gap=11)
    rect(s, Inches(7.5), Inches(1.9), Inches(5.2), Inches(4.1), DARK,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(7.85), Inches(2.2), Inches(4.6), Inches(3.6), [
        [("Marco normativo", 16, True, AMBER, HFONT)],
        [("", 6, False, WHITE)],
        [("Ley 19.628 — datos personales", 14, False, WHITE, BFONT)],
        [("Normativa CMF — ciberseguridad", 14, False, WHITE, BFONT)],
        [("OWASP Top 10 for LLM", 14, False, WHITE, BFONT)],
        [("", 8, False, WHITE)],
        [("Ética: transparencia, no maleficencia,", 13, False, ICE, BFONT, True)],
        [("equidad y supervisión humana.", 13, False, ICE, BFONT, True)],
    ])
    notas(s, "Recorre los 5 controles y el marco normativo. Reflexión ética en tus palabras: "
             "privacidad, confirmación humana, riesgo de sesgo. Prepárate para preguntas de "
             "dilemas éticos. (~1 min)")

    # ── 10. Recomendaciones ───────────────────────────────────────────────────
    s = slide_blank(prs)
    kicker(s, "9", "Mejoras propuestas")
    recs = [
        ("Desambiguación de intención", "clarificar + few-shot de enrutamiento (mayor precisión)"),
        ("Optimizar multi-agente", "paralelizar y cachear (menor latencia y costo)"),
        ("Resiliencia", "reintentos con backoff + modelo de respaldo"),
        ("Reforzar guardrail", "clasificador LLM + canary tokens (bloqueo →100%)"),
        ("Monitoreo proactivo", "alertas por umbrales sobre el dashboard"),
    ]
    y = Inches(1.8)
    for i, (t, d) in enumerate(recs, 1):
        rect(s, Inches(0.6), y, Inches(0.5), Inches(0.5), AMBER, shape=MSO_SHAPE.OVAL)
        text(s, Inches(0.6), y, Inches(0.5), Inches(0.5), [[(str(i), 16, True, DARK, HFONT)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.3), Emu(y - Inches(0.03)), Inches(11.2), Inches(0.95),
             [[(t + ".  ", 16, True, NAVY, BFONT), (d, 15, False, DARKTXT, BFONT)]],
             anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(y + Inches(1.02))
    notas(s, "Cada mejora se fundamenta en un hallazgo de las métricas. Prioriza la "
             "desambiguación (mayor impacto en precisión). Conecta con sostenibilidad y "
             "escalabilidad. (~1 min)")

    # ── 11. Cierre (dark) ─────────────────────────────────────────────────────
    s = slide_blank(prs)
    rect(s, 0, 0, W, H, DARK)
    rect(s, Inches(0.6), Inches(2.2), Inches(0.9), Inches(0.12), AMBER)
    text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(2.2), [
        [("Una solución medible, segura y escalable", 34, True, WHITE, HFONT)],
        [("LLM + RAG + agentes + observabilidad, de extremo a extremo.", 18, False, ICE, BFONT)],
    ])
    text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.5), [
        [("Repositorio: github.com/rogrosas/Ingenier-a-de-Soluciones-con-Inteligencia-Artificial",
          13, False, LGREY, BFONT)],
        [("", 6, False, WHITE)],
        [("Gracias — ¿preguntas?", 22, True, AMBER, HFONT)],
    ])
    notas(s, "Cierra con la idea fuerza: solución medible, segura y escalable. Invita a "
             "preguntas. Responde con evidencia y lenguaje técnico (IE13–IE15). (~30 s)")

    prs.save(SALIDA)
    return SALIDA, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    ruta, n = build()
    print(f"✔ Presentación generada: {os.path.relpath(ruta, ROOT)} ({n} diapositivas)")
